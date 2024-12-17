import os
import yfinance as yf
import matplotlib.pyplot as plt
import pandas as pd
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def get_font_with_fallback(preferred_font, is_title=False):
    """Get font name with fallback options"""
    try:
        # List of installed fonts - you might need to modify this based on your system
        available_fonts = ['Montserrat', 'Barlow', 'Arial', 'Calibri', 'Helvetica', 'Times New Roman']
        
        if preferred_font in available_fonts:
            return preferred_font
        else:
            # Default fallbacks
            if is_title:
                return 'Arial'  # Fallback for titles
            return 'Calibri'    # Fallback for body text
    except:
        return 'Arial' if is_title else 'Calibri'

def create_gradient_textbox(slide, left, top, width, height, text):
    """Create a textbox with gradient text and Barlow font with fallback"""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    run = p.add_run()
    run.text = text
    run.font.color.rgb = RGBColor(171, 146, 255)  # Light purple-blue color
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.name = get_font_with_fallback('Barlow', is_title=True)
    
    return text_box

def apply_white_text_to_shape(shape):
    """Apply white color and Montserrat font with fallback to all text"""
    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                # Only apply white color to non-title text
                if run.font.size != Pt(32):  # If not title size
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.name = get_font_with_fallback('Montserrat', is_title=False)
    elif hasattr(shape, "table"):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.font.name = get_font_with_fallback('Montserrat', is_title=False)
                        
def apply_slide_background(slide):
    """Apply dark background color to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(44, 48, 53)  # #2C3035

def apply_white_text_to_slide(slide):
    """Apply white text color to all shapes in a slide"""
    for shape in slide.shapes:
        if shape != slide.shapes.title:  # Skip the title shape if it exists
            apply_white_text_to_shape(shape)

class StockPresentationGenerator:
    def __init__(self, ticker):
        self.ticker = ticker.upper()
        try:
            self.stock = yf.Ticker(self.ticker)
            self.company_info = self.get_company_info()
        except Exception as e:
            print(f"Error initializing stock data: {e}")
            self.company_info = {'name': self.ticker}

    def get_company_info(self):
        """Retrieve comprehensive company information"""
        try:
            info = self.stock.info
            return {
                'name': info.get('longName', self.ticker),
                'industry': info.get('industry', 'N/A'),
                'sector': info.get('sector', 'N/A'),
                'description': info.get('longBusinessSummary', 'No description available'),
                'market_cap': info.get('marketCap', 'N/A'),
                'shares_outstanding': info.get('sharesOutstanding', 'N/A'),
                'float_shares': info.get('floatShares', 'N/A')
            }
        except Exception as e:
            print(f"Error retrieving company info: {e}")
            return {'name': self.ticker}

    def create_stock_price_chart(self):
        """Create stock price chart with moving averages"""
        try:
            hist = self.stock.history(period="1y")

            plt.style.use('dark_background')
            fig = plt.figure(figsize=(10, 6))
            
            # Set the figure and axes background color
            ax = plt.gca()
            fig.patch.set_facecolor('#2C3035')
            ax.set_facecolor('#2C3035')
            
            # Plot lines with custom colors
            plt.plot(hist.index, hist['Close'], label='Closing Price', color='#00A1C2', linewidth=2)
            plt.plot(hist.index, hist['Close'].rolling(window=50).mean(), 
                    label='50-day MA', color='#7F00B5', linewidth=1.5)
            plt.plot(hist.index, hist['Close'].rolling(window=200).mean(), 
                    label='200-day MA', color='#000086', linewidth=1.5)

            # Style the plot
            plt.title(f"{self.ticker} Stock Price", fontweight='bold', color='white', size=14)
            plt.xlabel("Date", color='white', size=12)
            plt.ylabel("Price", color='white', size=12)
            plt.grid(True, alpha=0.2)
            
            # Style the legend
            plt.legend(facecolor='#2C3035', edgecolor='white')
            
            # Style the axis
            ax.tick_params(colors='white')
            for spine in ax.spines.values():
                spine.set_color('white')

            # Save the chart
            chart_filename = f"{self.ticker}_stock_chart.png"
            plt.savefig(chart_filename, 
                       dpi=300, 
                       bbox_inches='tight', 
                       facecolor='#2C3035',
                       edgecolor='none')
            plt.close()

            return chart_filename
        except Exception as e:
            print(f"Error creating stock price chart: {e}")
            return None

    def get_financial_data(self):
            """Retrieve financial data for the past 4 years"""
            try:
                # Get financial data
                financials = self.stock.financials
                balance_sheet = self.stock.balance_sheet
                
                # Check if we have any data
                if financials.empty or balance_sheet.empty:
                    print("No financial data available")
                    return pd.DataFrame()

                # Get available years (up to 4)
                years = sorted(financials.columns, reverse=True)
                years = years[:min(4, len(years))]  # Take up to 4 years
                
                if not years:
                    print("No yearly data available")
                    return pd.DataFrame()

                data = []
                for year in years:
                    try:
                        # Safely get financial metrics
                        revenue = financials.loc['Total Revenue', year] / 1e9 if 'Total Revenue' in financials.index else 0
                        ebit = financials.loc['Operating Income', year] / 1e9 if 'Operating Income' in financials.index else 0
                        net_profit = financials.loc['Net Income', year] / 1e9 if 'Net Income' in financials.index else 0
                        
                        # Safely get total assets
                        total_assets = balance_sheet.loc['Total Assets', year] / 1e9 if 'Total Assets' in balance_sheet.index else 0
                        
                        # Calculate ROI
                        roi = (net_profit / total_assets * 100) if total_assets != 0 else 0

                        row = {
                            'Year': year.year,
                            'Revenue': f"${revenue:.1f}B" if revenue != 0 else "N/A",
                            'EBIT': f"${ebit:.1f}B" if ebit != 0 else "N/A",
                            'Net Profit': f"${net_profit:.1f}B" if net_profit != 0 else "N/A",
                            'ROI': f"{roi:.1f}%" if roi != 0 else "N/A"
                        }
                        data.append(row)
                    except Exception as e:
                        print(f"Warning: Error processing financial data for {year}: {e}")
                        row = {
                            'Year': year.year if hasattr(year, 'year') else 'N/A',
                            'Revenue': 'N/A',
                            'EBIT': 'N/A',
                            'Net Profit': 'N/A',
                            'ROI': 'N/A'
                        }
                        data.append(row)

                return pd.DataFrame(data)
            except Exception as e:
                print(f"Error retrieving financial data: {e}")
                return pd.DataFrame()
            
    def get_analyst_insights(self):
        """Retrieve analyst insights"""
        try:
            info = self.stock.info
            current_price = info.get('currentPrice', 'N/A')
            target_price = info.get('targetMeanPrice', 'N/A')

            upside = 'N/A'
            if current_price != 'N/A' and target_price != 'N/A':
                try:
                    upside = ((target_price - current_price) / current_price) * 100
                except:
                    pass

            return {
                'rating_numeric': info.get('recommendationMean', 'N/A'),
                'recommendation': self._translate_recommendation(info.get('recommendationKey', 'N/A')),
                'number_of_analysts': info.get('numberOfAnalystOpinions', 'N/A'),
                'mean_target_price': target_price,
                'current_price': current_price,
                'implied_upside': f"{upside:.1f}%" if isinstance(upside, float) else upside
            }
        except Exception as e:
            print(f"Error retrieving analyst insights: {e}")
            return {
                'rating_numeric': 'N/A',
                'recommendation': 'N/A',
                'number_of_analysts': 'N/A',
                'mean_target_price': 'N/A',
                'current_price': 'N/A',
                'implied_upside': 'N/A'
            }

    def _translate_recommendation(self, rec_key):
        """Translate recommendation key"""
        translations = {
            'buy': 'Buy',
            'hold': 'Hold',
            'sell': 'Sell',
            'strongBuy': 'Strong Buy',
            'strongSell': 'Strong Sell'
        }
        return translations.get(str(rec_key).lower(), rec_key)

    def get_financial_health(self):
        """Retrieve financial health metrics"""
        try:
            info = self.stock.info
            return {
                'cash_position': info.get('totalCash', 'N/A'),
                'total_debt': info.get('totalDebt', 'N/A'),
                'current_ratio': info.get('currentRatio', 'N/A')
            }
        except Exception as e:
            print(f"Error retrieving financial health: {e}")
            return {
                'cash_position': 'N/A',
                'total_debt': 'N/A',
                'current_ratio': 'N/A'
            }

    def get_dividend_info(self):
        """Retrieve dividend information"""
        try:
            info = self.stock.info
            ex_div_date = info.get('exDividendDate', 'N/A')
            if isinstance(ex_div_date, (int, float)):
                ex_div_date = datetime.fromtimestamp(ex_div_date).strftime('%Y-%m-%d')

            div_rate = info.get('dividendRate', 'N/A')
            div_yield = info.get('dividendYield', 'N/A')

            return {
                'ex_dividend_date': ex_div_date,
                'dividend_rate': f"${div_rate:.2f}" if isinstance(div_rate, (int, float)) else 'N/A',
                'dividend_yield': f"{div_yield*100:.2f}%" if isinstance(div_yield, (int, float)) else 'N/A'
            }
        except Exception as e:
            print(f"Error retrieving dividend information: {e}")
            return {
                'ex_dividend_date': 'N/A',
                'dividend_rate': 'N/A',
                'dividend_yield': 'N/A'
            }

    def generate_key_strengths_and_catalysts(self):
        """Generate strengths and catalysts"""
        try:
            industry = self.company_info.get('industry', 'the industry')
            return {
                'strengths': [
                    f"Strong market position in {industry}",
                    "Innovative product development strategy",
                    "Robust financial performance"
                ],
                'catalysts': [
                    "Potential expansion into emerging markets",
                    "Technological innovation pipeline",
                    "Strategic partnerships and acquisitions"
                ]
            }
        except Exception as e:
            print(f"Error generating strengths and catalysts: {e}")
            return {
                'strengths': ["Data not available"],
                'catalysts': ["Data not available"]
            }

    def generate_risk_analysis(self):
        """Generate risk analysis"""
        try:
            return {
                'risks': [
                    "Market volatility and economic uncertainties",
                    "Intense competitive landscape",
                    "Regulatory challenges"
                ],
                'mitigations': [
                    "Diversified revenue streams",
                    "Strong R&D investment",
                    "Robust compliance framework"
                ]
            }
        except Exception as e:
            print(f"Error generating risk analysis: {e}")
            return {
                'risks': ["Data not available"],
                'mitigations': ["Data not available"]
            }

    def generate_recommendation(self, strengths, catalysts):
        """Generate recommendation"""
        try:
            analyst_insights = self.get_analyst_insights()
            recommendation = analyst_insights['recommendation']

            return f"""Based on our comprehensive analysis, we recommend a {recommendation} rating for {self.company_info['name']}.

Key Highlights:
• {strengths[0]}
• {strengths[1]}
• {catalysts[0]}

The company shows promising potential with strong market positioning and innovative strategies."""
        except Exception as e:
            print(f"Error generating recommendation: {e}")
            return "Unable to generate recommendation due to insufficient data"

    def _add_title_to_slide(self, slide, title_text):
        """Add styled title to slide"""
        apply_slide_background(slide)
        
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        
        return create_gradient_textbox(slide, left, top, width, height, title_text)

 
    def _generate_overview_slide(self, prs):
            """Generate overview slide"""
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_title_to_slide(slide, "Company Overview")

            # Basic info box
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(3.5)
            height = Inches(2)

            basic_info_box = slide.shapes.add_textbox(left, top, width, height)
            basic_tf = basic_info_box.text_frame

            # Company info with bold labels
            p1 = basic_tf.add_paragraph()
            run1 = p1.add_run()
            run1.text = "Company: "
            run1.font.bold = True
            run2 = p1.add_run()
            run2.text = self.company_info['name']
            p1.add_line_break()
            p1.add_line_break()

            # Industry
            p2 = basic_tf.add_paragraph()
            run3 = p2.add_run()
            run3.text = "Industry: "
            run3.font.bold = True
            run4 = p2.add_run()
            run4.text = self.company_info['industry']
            p2.add_line_break()
            p2.add_line_break()

            # Sector
            p3 = basic_tf.add_paragraph()
            run5 = p3.add_run()
            run5.text = "Sector: "
            run5.font.bold = True
            run6 = p3.add_run()
            run6.text = self.company_info['sector']

            # Stock price chart
            chart_path = self.create_stock_price_chart()
            if chart_path:
                slide.shapes.add_picture(chart_path, Inches(4.5), Inches(1.5), width=Inches(5))
                os.remove(chart_path)

            # Description box
            desc_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.75), Inches(9), Inches(1.5))
            desc_tf = desc_box.text_frame

            # Add description
            description = self.company_info['description']
            sentences = description.split('.')
            limited_description = '. '.join(sentence.strip() 
                                        for sentence in sentences[:3] if sentence.strip()) + '.'

            para = desc_tf.add_paragraph()
            run_desc1 = para.add_run()
            run_desc1.text = "Description: "
            run_desc1.font.bold = True
            run_desc2 = para.add_run()
            run_desc2.text = limited_description
            desc_tf.word_wrap = True

            apply_white_text_to_slide(slide)

    def _generate_market_position_slide(self, prs):
        """Generate market position slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title_to_slide(slide, "Market Position")

        # Content box
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(6), Inches(5))
        tf = content_box.text_frame

        # Format values
        market_cap_t = self.company_info['market_cap'] / 1e12 if isinstance(
            self.company_info['market_cap'], (int, float)) else 'N/A'
        shares_out_b = self.company_info['shares_outstanding'] / 1e9 if isinstance(
            self.company_info['shares_outstanding'], (int, float)) else 'N/A'
        float_shares_b = self.company_info['float_shares'] / 1e9 if isinstance(
            self.company_info['float_shares'], (int, float)) else 'N/A'

        # Market Statistics header
        p1 = tf.add_paragraph()
        run1 = p1.add_run()
        run1.text = "Market Statistics"
        run1.font.bold = True
        run1.font.underline = True
        run1.font.size = Pt(18)
        p1.add_line_break()
        p1.add_line_break()

        # Market cap
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = "Market Cap: "
        run2.font.bold = True
        run3 = p2.add_run()
        run3.text = f"${market_cap_t:.1f}T" if isinstance(market_cap_t, float) else "N/A"
        p2.add_line_break()
        p2.add_line_break()

        # Shares outstanding
        p3 = tf.add_paragraph()
        run4 = p3.add_run()
        run4.text = "Shares Outstanding: "
        run4.font.bold = True
        run5 = p3.add_run()
        run5.text = f"{shares_out_b:.1f}B" if isinstance(shares_out_b, float) else "N/A"
        p3.add_line_break()
        p3.add_line_break()

        # Float shares
        p4 = tf.add_paragraph()
        run6 = p4.add_run()
        run6.text = "Float Shares: "
        run6.font.bold = True
        run7 = p4.add_run()
        run7.text = f"{float_shares_b:.1f}B" if isinstance(float_shares_b, float) else "N/A"

        # Financial table
        financial_data = self.get_financial_data()
        if not financial_data.empty:
            table = slide.shapes.add_table(
                len(financial_data) + 1,
                len(financial_data.columns),
                Inches(4), Inches(2),
                Inches(5.6), Inches(4)
            ).table

            # Add headers
            for col, header in enumerate(financial_data.columns):
                cell = table.cell(0, col)
                cell.text = header
                cell.text_frame.paragraphs[0].runs[0].font.bold = True
                if header == 'Revenue':
                    table.columns[col].width = table.columns[col].width + Inches(0.1)

            # Add data
            for row, (_, data_row) in enumerate(financial_data.iterrows(), start=1):
                for col, value in enumerate(data_row):
                    table.cell(row, col).text = str(value)

        apply_white_text_to_slide(slide)

    def _generate_strengths_catalysts_slide(self, prs):
        """Generate strengths and catalysts slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title_to_slide(slide, "Key Strengths and Growth Catalysts")

        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf = content_box.text_frame

        strengths_catalysts = self.generate_key_strengths_and_catalysts()

        # Key Strengths section
        p1 = tf.add_paragraph()
        run1 = p1.add_run()
        run1.text = "Key Strengths:"
        run1.font.bold = True
        p1.add_line_break()

        for strength in strengths_catalysts['strengths']:
            p = tf.add_paragraph()
            p.text = f"• {strength}"

        tf.add_paragraph()  # Empty line

        # Growth Catalysts section
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = "Growth Catalysts:"
        run2.font.bold = True
        p2.add_line_break()

        for catalyst in strengths_catalysts['catalysts']:
            p = tf.add_paragraph()
            p.text = f"• {catalyst}"

        apply_white_text_to_slide(slide)

    def _generate_investment_thesis_slide(self, prs):
        """Generate investment thesis slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title_to_slide(slide, "Investment Thesis")

        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.05), Inches(9), Inches(5))
        tf = content_box.text_frame

        # Get and format all data
        financial_health = self.get_financial_health()
        analyst_insights = self.get_analyst_insights()
        dividend_info = self.get_dividend_info()

        # Financial Health section
        p1 = tf.add_paragraph()
        run1 = p1.add_run()
        run1.text = "Financial Health:"
        run1.font.bold = True
        tf.add_paragraph().font.size = Pt(9)  # Spacing

        for key, value in financial_health.items():
            p = tf.add_paragraph()
            p.text = f"{key.replace('_', ' ').title()}: {value}"

        tf.add_paragraph().font.size = Pt(9)  # Spacing

        # Analyst Insights section
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = "Analyst Insights:"
        run2.font.bold = True
        tf.add_paragraph().font.size = Pt(9)  # Spacing

        for key, value in analyst_insights.items():
            p = tf.add_paragraph()
            p.text = f"{key.replace('_', ' ').title()}: {value}"

        tf.add_paragraph().font.size = Pt(9)  # Spacing

        # Dividend Information section
        p3 = tf.add_paragraph()
        run3 = p3.add_run()
        run3.text = "Dividend Information:"
        run3.font.bold = True
        tf.add_paragraph().font.size = Pt(9)  # Spacing

        for key, value in dividend_info.items():
            p = tf.add_paragraph()
            p.text = f"{key.replace('_', ' ').title()}: {value}"

        apply_white_text_to_slide(slide)

    def _generate_risk_analysis_slide(self, prs):
        """Generate risk analysis slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title_to_slide(slide, "Risk Analysis and Mitigation")

        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf = content_box.text_frame

        risk_analysis = self.generate_risk_analysis()

        # Risks section
        p1 = tf.add_paragraph()
        run1 = p1.add_run()
        run1.text = "Risks:"
        run1.font.bold = True
        p1.add_line_break()

        for risk in risk_analysis['risks']:
            p = tf.add_paragraph()
            p.text = f"• {risk}"

        tf.add_paragraph()  # Spacing

        # Mitigation section
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = "Mitigation Strategies:"
        run2.font.bold = True
        p2.add_line_break()

        for mitigation in risk_analysis['mitigations']:
            p = tf.add_paragraph()
            p.text = f"• {mitigation}"

        apply_white_text_to_slide(slide)

    def _generate_conclusion_slide(self, prs):
        """Generate conclusion slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title_to_slide(slide, "Conclusion and Recommendation")

        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf = content_box.text_frame

        strengths_catalysts = self.generate_key_strengths_and_catalysts()
        recommendation = self.generate_recommendation(
            strengths_catalysts['strengths'],
            strengths_catalysts['catalysts']
        )

        p = tf.add_paragraph()
        p.text = recommendation
        
        apply_white_text_to_slide(slide)

    def generate_presentation(self):
            prs = Presentation()

            # Title slide
            slide1 = prs.slides.add_slide(prs.slide_layouts[6])
            apply_slide_background(slide1)

            # Create centered title with gradient
            left = Inches(1)
            top = Inches(3)
            width = Inches(8)
            height = Inches(1.5)

            title_text = f"{self.company_info.get('name', self.ticker)} Stock Analysis"
            create_gradient_textbox(slide1, left, top, width, height, title_text)

            # Generate content slides - with error handling for each
            try:
                self._generate_overview_slide(prs)
            except Exception as e:
                print(f"Warning: Error generating overview slide: {e}")

            try:
                self._generate_market_position_slide(prs)
            except Exception as e:
                print(f"Warning: Error generating market position slide: {e}")

            try:
                self._generate_strengths_catalysts_slide(prs)
            except Exception as e:
                print(f"Warning: Error generating strengths slide: {e}")

            try:
                self._generate_investment_thesis_slide(prs)
            except Exception as e:
                print(f"Warning: Error generating investment thesis slide: {e}")

            try:
                self._generate_risk_analysis_slide(prs)
            except Exception as e:
                print(f"Warning: Error generating risk analysis slide: {e}")

            try:
                self._generate_conclusion_slide(prs)
            except Exception as e:
                print(f"Warning: Error generating conclusion slide: {e}")

            # Save presentation
            output_filename = f"{self.ticker}_stock_analysis.pptx"
            prs.save(output_filename)
            return output_filename


def main():
    """Main function to run the presentation generator"""
    ticker = input("Enter stock ticker: ")
    generator = StockPresentationGenerator(ticker)
    
    presentation_path = generator.generate_presentation()
    if presentation_path:
        print(f"Presentation generated successfully: {presentation_path}")
    else:
        print("Failed to generate presentation")
    

if __name__ == "__main__":
    main()